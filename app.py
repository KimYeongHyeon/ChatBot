import streamlit as st
import base64
import os
from io import BytesIO
import PyPDF2
import docx
import pptx
import requests

# Function to read and process the uploaded file
def process_file(file, file_type):
    if file_type == "pdf":
        pdf = PyPDF2.PdfFileReader(file)
        text = ""
        for page in range(pdf.getNumPages()):
            text += pdf.getPage(page).extractText()
    elif file_type == "docx":
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_type == "pptx":
        ppt = pptx.Presentation(file)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    elif file_type == "txt":
        text = file.read().decode("utf-8")
    return text

# Streamlit app
# Streamlit app
st.title("File Upload and Chat Application")

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if 'file_text' not in st.session_state:
    file = st.file_uploader("Upload a PDF, Word, PPT, or Text file", type=["pdf", "docx", "pptx", "txt"])

    if file:
        file_type = file.type.split("/")[-1]
        text = process_file(file, file_type)
        st.session_state.file_text = text
        st.experimental_rerun()
else:
    st.write("File upload completed.")

st.header("Chat with API")
chat_input = st.text_input("Type your message here")

if st.button("Send"):
    # response = requests.post("https://your-chat-api-url.com", json={"message": chat_input})
    # chat_response = response.json().get("response", "Sorry, I couldn't understand your message.")
    chat_response = chat_input
    st.session_state.chat_history.append({"user": chat_input, "hermes": chat_response})

for chat in st.session_state.chat_history:
    st.write("You:", chat["user"])
    st.write("Hermes:", chat["hermes"])
